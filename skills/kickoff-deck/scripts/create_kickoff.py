#!/usr/bin/env python3
"""
Generate a Project+ kickoff deck .pptx from a JSON input.

Builds a 9-10 slide deck: title, vision, objectives, scope, team,
milestones, risks, ways of working, asks, plus a Draft Gaps slide
at position 2 when the draft has unresolved items.

Uses python-pptx for consistency with the other document skills.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


TBD_PATTERN = re.compile(r"\[TBD[^\]]*\]")

# Figure-lint patterns — any token matching one of these in the draft must
# also appear literally in the --source text, or it is flagged unverified.
FIGURE_PATTERNS = [
    re.compile(r"\$\s?\d+(?:[,\.]\d+)*\s?[kKmMbB]?"),             # $150k, $1.2M, $150,000
    re.compile(r"\d+(?:\.\d+)?\s?%"),                              # 99.9%, 50 %
    re.compile(r"\b\d{4}-\d{2}-\d{2}\b"),                          # 2026-06-30
    re.compile(r"\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2}(?:,\s*\d{4})?\b"),  # June 30, 2026
    re.compile(r"\b\d{1,2}/\d{1,2}(?:/\d{2,4})?\b"),               # 6/30, 06/30/2026
]


def lint_figures_against_source(data: dict, source_text: str) -> list[str]:
    """Flag any figure in the draft that isn't present verbatim in source."""
    if not source_text:
        return []
    unverified: list[str] = []
    seen: set[str] = set()

    def walk(node, path="root"):
        if isinstance(node, dict):
            for k, v in node.items():
                walk(v, f"{path}.{k}")
        elif isinstance(node, list):
            for i, v in enumerate(node):
                walk(v, f"{path}[{i}]")
        elif isinstance(node, str):
            if TBD_PATTERN.search(node):
                return
            # Skip metadata fields that are not user-facing figures.
            if path in ("root.header.date", "root.header.version"):
                return
            for pat in FIGURE_PATTERNS:
                for m in pat.findall(node):
                    token = m.strip()
                    if not token or token in seen:
                        continue
                    if token not in source_text:
                        seen.add(token)
                        unverified.append(f"⚠ unverified figure '{token}' in {path} — not found in source")

    walk(data)
    return unverified

TBD_COLOR = RGBColor(0xB7, 0x47, 0x1E)      # burnt orange
HEADING_COLOR = RGBColor(0x1F, 0x3A, 0x5F)  # deep navy
BODY_COLOR = RGBColor(0x33, 0x33, 0x33)     # near-black
MUTED_COLOR = RGBColor(0x66, 0x66, 0x66)    # gray


def is_tbd(value) -> bool:
    if value is None:
        return True
    if isinstance(value, str):
        return bool(TBD_PATTERN.search(value)) or not value.strip()
    if isinstance(value, (list, dict)):
        return len(value) == 0
    return False


def collect_tbds(data: dict) -> list[str]:
    gaps: list[str] = []
    header = data.get("header", {}) or {}
    for k, v in header.items():
        if is_tbd(v):
            gaps.append(f"header.{k}: {v or '[TBD — empty]'}")
    for key in ("vision",):
        if is_tbd(data.get(key)):
            gaps.append(f"{key}: {data.get(key) or '[TBD — empty]'}")
    for key in ("objectives", "asks"):
        arr = data.get(key) or []
        if not arr:
            gaps.append(f"{key}: (empty)")
        for i, v in enumerate(arr):
            if is_tbd(v):
                gaps.append(f"{key}[{i}]: {v}")
    scope = data.get("scope", {}) or {}
    for sub in ("in_scope", "out_of_scope"):
        arr = scope.get(sub) or []
        if not arr:
            gaps.append(f"scope.{sub}: (empty)")
        for i, v in enumerate(arr):
            if is_tbd(v):
                gaps.append(f"scope.{sub}[{i}]: {v}")
    team = data.get("team") or []
    if not team:
        gaps.append("team: (empty)")
    for i, m in enumerate(team):
        for f in ("name", "role"):
            if is_tbd(m.get(f)):
                gaps.append(f"team[{i}].{f}: {m.get(f)}")
    milestones = data.get("milestones") or []
    if not milestones:
        gaps.append("milestones: (empty)")
    for i, m in enumerate(milestones):
        for f in ("name", "target"):
            if is_tbd(m.get(f)):
                gaps.append(f"milestones[{i}].{f}: {m.get(f)}")
    risks = data.get("risks") or []
    if not risks:
        gaps.append("risks: (empty)")
    for i, r in enumerate(risks):
        for f in ("risk", "mitigation"):
            if is_tbd(r.get(f)):
                gaps.append(f"risks[{i}].{f}: {r.get(f)}")
    wow = data.get("ways_of_working", {}) or {}
    for f in ("cadence", "tools", "decision_forum"):
        if is_tbd(wow.get(f)):
            gaps.append(f"ways_of_working.{f}: {wow.get(f)}")
    return gaps


# ---------- slide builders ----------


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN * 2


def _add_text(slide, left, top, width, height, text, *,
              size=18, bold=False, color=BODY_COLOR, align=PP_ALIGN.LEFT,
              italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic or (isinstance(text, str) and bool(TBD_PATTERN.search(text)))
    if isinstance(text, str) and TBD_PATTERN.search(text):
        run.font.color.rgb = TBD_COLOR
    else:
        run.font.color.rgb = color
    return box


def _responsive_size(n_items: int, base_size: int) -> tuple[int, int]:
    """Return (font_size, space_after_pt) scaled to item count.

    Sparse slides (1–2 items) get much larger text so a short list doesn't
    leave 4 inches of whitespace. Dense slides (5+) use the base size.
    Floor: 1 item → 44pt, 2 items → 36pt, 3 items → 28pt.
    """
    if n_items <= 1:
        return (max(base_size + 24, 44), 18)
    if n_items == 2:
        return (max(base_size + 16, 36), 14)
    if n_items == 3:
        return (max(base_size + 8, 28), 10)
    if n_items == 4:
        return (base_size + 3, 8)
    return (base_size, 6)


def _add_bullets(slide, left, top, width, height, items, *,
                 size=18, color=BODY_COLOR, anchor_middle=True, responsive=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor_middle:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    items = items or ["[TBD — empty]"]
    if responsive:
        font_size, space_after = _responsive_size(len(items), size)
    else:
        font_size, space_after = (size, 6)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•  {item}" if item is not None else "•  [TBD — empty]"
        run.font.size = Pt(font_size)
        if isinstance(item, str) and TBD_PATTERN.search(item):
            run.font.color.rgb = TBD_COLOR
            run.font.italic = True
        else:
            run.font.color.rgb = color
        p.space_after = Pt(space_after)


def _slide_header(slide, title, subtitle=None):
    _add_text(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.7),
              title, size=32, bold=True, color=HEADING_COLOR)
    if subtitle:
        _add_text(slide, MARGIN, Inches(1.1), CONTENT_W, Inches(0.4),
                  subtitle, size=14, italic=True, color=MUTED_COLOR)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_title_slide(prs, header):
    slide = _blank_slide(prs)
    name = header.get("project_name") or "[TBD — no project name]"
    _add_text(slide, MARGIN, Inches(2.5), CONTENT_W, Inches(1.2),
              name, size=48, bold=True, color=HEADING_COLOR, align=PP_ALIGN.CENTER)
    _add_text(slide, MARGIN, Inches(3.8), CONTENT_W, Inches(0.6),
              "Project Kickoff", size=24, italic=True, color=MUTED_COLOR,
              align=PP_ALIGN.CENTER)
    meta_line = (
        f"PM: {header.get('project_manager', '[TBD]')}    "
        f"Sponsor: {header.get('sponsor', '[TBD]')}    "
        f"Date: {header.get('date', '[TBD]')}"
    )
    _add_text(slide, MARGIN, Inches(5.0), CONTENT_W, Inches(0.5),
              meta_line, size=14, color=BODY_COLOR, align=PP_ALIGN.CENTER)


def add_gaps_slide(prs, gaps):
    slide = _blank_slide(prs)
    _slide_header(slide, f"Draft Gaps ({len(gaps)})",
                  "Unresolved items — resolve before using this deck with stakeholders.")
    # Show up to 14 gaps; overflow goes to open_questions.md
    shown = gaps[:14]
    _add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.5),
                 shown, size=14, color=TBD_COLOR,
                 anchor_middle=False, responsive=False)
    if len(gaps) > 14:
        _add_text(slide, MARGIN, Inches(7.0), CONTENT_W, Inches(0.4),
                  f"+ {len(gaps) - 14} more — see open_questions.md",
                  size=12, italic=True, color=MUTED_COLOR)


def add_vision_slide(prs, vision):
    slide = _blank_slide(prs)
    _slide_header(slide, "Why we're here")
    _add_text(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(3.0),
              vision or "[TBD — vision not stated]", size=24, color=BODY_COLOR)


def add_objectives_slide(prs, objectives):
    slide = _blank_slide(prs)
    _slide_header(slide, "Objectives & success")
    _add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.5),
                 objectives, size=20)


def add_scope_slide(prs, scope):
    slide = _blank_slide(prs)
    _slide_header(slide, "Scope")
    col_w = (CONTENT_W - Inches(0.4)) / 2
    _add_text(slide, MARGIN, Inches(1.7), col_w, Inches(0.5),
              "In scope", size=18, bold=True, color=HEADING_COLOR)
    _add_bullets(slide, MARGIN, Inches(2.2), col_w, Inches(5.0),
                 scope.get("in_scope") or [], size=16)
    left2 = MARGIN + col_w + Inches(0.4)
    _add_text(slide, left2, Inches(1.7), col_w, Inches(0.5),
              "Out of scope", size=18, bold=True, color=HEADING_COLOR)
    _add_bullets(slide, left2, Inches(2.2), col_w, Inches(5.0),
                 scope.get("out_of_scope") or [], size=16)


def add_team_slide(prs, team):
    slide = _blank_slide(prs)
    _slide_header(slide, "Team")
    lines = []
    if not team:
        lines = ["[TBD — team not named]"]
    else:
        for m in team:
            name = m.get("name", "[TBD]")
            role = m.get("role", "[TBD]")
            lines.append(f"{name} — {role}")
    _add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.5),
                 lines, size=18)


def add_milestones_slide(prs, milestones):
    slide = _blank_slide(prs)
    _slide_header(slide, "Milestones")
    lines = []
    if not milestones:
        lines = ["[TBD — milestones not stated]"]
    else:
        for m in milestones:
            lines.append(f"{m.get('name', '[TBD]')}  —  {m.get('target', '[TBD]')}")
    _add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.5),
                 lines, size=18)


def add_risks_slide(prs, risks):
    slide = _blank_slide(prs)
    _slide_header(slide, "Risks & mitigations")
    lines = []
    if not risks:
        lines = ["[TBD — risks not stated]"]
    else:
        for r in risks:
            lines.append(f"{r.get('risk', '[TBD]')}  →  {r.get('mitigation', '[TBD]')}")
    _add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.5),
                 lines, size=16)


def add_wow_slide(prs, wow):
    slide = _blank_slide(prs)
    _slide_header(slide, "Ways of working")
    lines = [
        f"Cadence: {wow.get('cadence', '[TBD]')}",
        f"Tools: {wow.get('tools', '[TBD]')}",
        f"Decision forum: {wow.get('decision_forum', '[TBD]')}",
    ]
    _add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.5),
                 lines, size=20)


def add_asks_slide(prs, asks):
    slide = _blank_slide(prs)
    _slide_header(slide, "What we need from you today")
    _add_bullets(slide, MARGIN, Inches(1.7), CONTENT_W, Inches(5.5),
                 asks, size=20)


def render_kickoff(data: dict, output_path: Path, source_text: str = "") -> list[str]:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    header = data.get("header", {}) or {}
    gaps = collect_tbds(data)
    if source_text:
        gaps.extend(lint_figures_against_source(data, source_text))

    add_title_slide(prs, header)
    if gaps:
        add_gaps_slide(prs, gaps)
    add_vision_slide(prs, data.get("vision"))
    add_objectives_slide(prs, data.get("objectives") or [])
    add_scope_slide(prs, data.get("scope", {}) or {})
    add_team_slide(prs, data.get("team") or [])
    add_milestones_slide(prs, data.get("milestones") or [])
    add_risks_slide(prs, data.get("risks") or [])
    add_wow_slide(prs, data.get("ways_of_working", {}) or {})
    add_asks_slide(prs, data.get("asks") or [])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return gaps


def write_open_questions(gaps: list[str], output_path: Path) -> None:
    lines = [
        "# Open Questions for Sponsor / PM",
        "",
        f"This kickoff deck draft has **{len(gaps)} unresolved items**. "
        "Resolve these before using the deck with stakeholders.",
        "",
    ]
    by_section: dict[str, list[str]] = {}
    for g in gaps:
        section = g.split(".")[0].split("[")[0].strip() or "general"
        by_section.setdefault(section, []).append(g)
    for section, items in sorted(by_section.items()):
        lines.append(f"## {section}")
        lines.append("")
        for item in items:
            lines.append(f"- {item}")
        lines.append("")
    output_path.write_text("\n".join(lines))


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate a project kickoff deck .pptx")
    parser.add_argument("--input", help="Path to JSON input (omit to read stdin)")
    parser.add_argument("--output", required=True)
    parser.add_argument("--questions-threshold", type=int, default=5)
    parser.add_argument("--source", help="Path to user's raw message text; enables figure linting")
    args = parser.parse_args()

    source_text = ""
    if args.source:
        try:
            source_text = Path(args.source).read_text()
        except Exception as e:
            print(f"WARN: could not read --source file: {e}", file=sys.stderr)

    output_path = Path(args.output)

    try:
        if args.input:
            data = json.loads(Path(args.input).read_text())
        else:
            data = json.loads(sys.stdin.read())
    except Exception as e:
        print(f"ERROR: could not parse input JSON: {e}", file=sys.stderr)
        return 2

    gaps = render_kickoff(data, output_path, source_text)
    print(f"Kickoff deck written: {output_path}")
    print(f"TBD count: {len(gaps)}")

    if len(gaps) >= args.questions_threshold:
        oq = output_path.parent / "open_questions.md"
        write_open_questions(gaps, oq)
        print(f"Open questions written: {oq}")

    return 0


if __name__ == "__main__":
    sys.exit(main())
